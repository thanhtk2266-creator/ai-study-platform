import os
import aiofiles
from fastapi import UploadFile
from sqlalchemy.orm import Session
from pathlib import Path
import uuid
import PyPDF2
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter

from src.core.config import settings
from src.core.database import SessionLocal
from src.models.document import Document
from src.services.ai_service import AIService

class DocumentService:
    def __init__(self, ai_service: AIService):
        self.ai_service = ai_service

    async def save_upload_file(self, file: UploadFile) -> str:
        os.makedirs(settings.UPLOAD_DIR, exist_ok=True)
        content = await file.read()
        if len(content) > settings.MAX_FILE_SIZE:
            raise ValueError("File exceeds max upload size")

        safe_name = Path(file.filename or "upload.bin").name
        file_path = os.path.join(settings.UPLOAD_DIR, f"{uuid.uuid4()}_{safe_name}")

        async with aiofiles.open(file_path, 'wb') as out_file:
            await out_file.write(content)
        return file_path

    def parse_document(self, file_path: str) -> str:
        text = ""
        ext = file_path.lower().split('.')[-1]
        
        if ext == "pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        elif ext in ["ppt", "pptx"]:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        
        return text

    def chunk_text(self, text: str) -> list[str]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )
        return splitter.split_text(text)

    def process_document(self, document_id: str, file_path: str):
        db: Session = SessionLocal()
        try:
            # 1. Trích xuất text
            text = self.parse_document(file_path)
            
            # 2. Chia nhỏ text
            chunks = self.chunk_text(text)
            
            # 3. Lưu vào ChromaDB
            self.ai_service.store_embeddings(document_id, chunks)
            
            # 4. Cập nhật trạng thái
            doc = db.query(Document).filter(Document.id == document_id).first()
            if doc:
                doc.status = "ready"
                doc.chunk_count = len(chunks)
                db.commit()
                
        except Exception as e:
            doc = db.query(Document).filter(Document.id == document_id).first()
            if doc:
                doc.status = "error"
                db.commit()
            print(f"Error processing document: {e}")
        finally:
            db.close()
